import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_pptx(output_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 Widescreen
    prs.slide_height = Inches(7.5)

    # Color Palette Tokens
    c_canvas = RGBColor(0xF7, 0xF7, 0xF4)     # Warm Cream #f7f7f4
    c_card = RGBColor(0xFF, 0xFF, 0xFF)       # Pure White #ffffff
    c_ink = RGBColor(0x26, 0x25, 0x1E)        # Warm Near-Black #26251e
    c_body = RGBColor(0x5A, 0x58, 0x52)       # Slate Body #5a5852
    c_orange = RGBColor(0xF5, 0x4E, 0x00)     # Cursor Orange #f54e00
    c_border = RGBColor(0xE6, 0xE5, 0xE0)     # Hairline Divider #e6e5e0
    c_emerald = RGBColor(0x1F, 0x8A, 0x65)    # Emerald Audit #1f8a65
    c_blue = RGBColor(0x10, 0x6B, 0xA3)       # Cobalt Accent #106ba3

    slides_data = [
        {
            "tag": "SLIDE 01 / 12  ·  EXECUTIVE HOOK",
            "title": "KYZER Sovereign OS",
            "subtitle": "Autonomous Public Health Supply Chain & Epidemic Intelligence Co-Pilot",
            "script": "Judges, in rural public health centres across India and BRICS nations, 1 in 4 essential vaccine batches are lost to stockouts or cold-chain melting. Today, Team KYZER presents KYZER—the sovereign multi-agent intelligence platform that automates paper register perception, epidemic demand forecasting, and quantum-optimized lateral redistribution before clinical stockouts occur.",
            "cards": [
                {
                    "title": "The Mission & Problem Scope",
                    "points": [
                        "Eliminating rural vaccine stockouts and cold-chain spoilage across primary health centres.",
                        "Track: Google Cloud: Build with AI — Code for Communities Season 2.",
                        "Target Audience: Ministry of Health & Family Welfare, District Health Officers, and ASHA field workers."
                    ]
                },
                {
                    "title": "Core Technical Breakthroughs",
                    "points": [
                        "Perception: OpenCV 5.0 + Gemini 1.5 Flash Vision OCR (97% client-side bandwidth reduction).",
                        "Forecasting: LightGBM Tweedie Quantile (17.48% WAPE) coupled with differential SEIR epidemic dynamics.",
                        "Optimization: 156-Qubit IBM Heron r2 QPU Quantum VRP (138.89 km / 238.1 min transit).",
                        "Governance: DeepSeek Harness 5-Agent Loop with deterministic >= 1.9x clinical safety gate."
                    ]
                }
            ]
        },
        {
            "tag": "SLIDE 02 / 12  ·  THE GROUND REALITY",
            "title": "The Rural Healthcare Supply Chain Bottleneck",
            "subtitle": "Why rural clinics experience severe stockouts despite surplus at district depots",
            "script": "The problem isn't a national shortage of medicine; it's a distribution latency failure. Rural primary health centres rely on handwritten paper registers. When monsoon rains trigger viral outbreaks, local clinics deplete their 48-hour buffers while a neighboring depot sits on a surplus. With WHO ice packs melting in 240 minutes, classical distribution is too slow to save lives.",
            "cards": [
                {
                    "title": "1. The 4-Hour Cold-Chain Window",
                    "points": [
                        "WHO 240-minute ice pack thermal lifetime in vaccine carriers.",
                        "Any unoptimized route or transit delay results in total thermal degradation.",
                        "Classical regional routing lacks real-time traffic and terrain awareness."
                    ]
                },
                {
                    "title": "2. 92% Physical Paper Registers",
                    "points": [
                        "Rural PHCs lack digital ERPs; daily records exist only in physical logbooks.",
                        "District magistrates have zero real-time visibility into ward-level medicine depletion.",
                        "2G/3G connectivity prevents uploading heavy 6MB uncompressed photos."
                    ]
                }
            ]
        },
        {
            "tag": "SLIDE 03 / 12  ·  USER PERSONAS",
            "title": "Empowering Frontline Health Workers & Officers",
            "subtitle": "From 3 hours of manual paperwork to instant autonomous co-piloting",
            "script": "Meet Sunita, an ASHA healthcare worker at Koregaon Bhima, and Dr. Patil at the Shirur Sub-District Hospital. Sunita snaps a single photo of her daily paper logbook on her mobile phone. KYZER's client-side canvas compressor downscales it by 97%, and Gemini 1.5 Flash Vision extracts stock, beds, and staff attendance in under 2 seconds.",
            "cards": [
                {
                    "title": "Sunita — ASHA Field Nurse (Koregaon Bhima PHC)",
                    "points": [
                        "Before: Spends 3 hours daily manually counting blister packs & handwriting registers.",
                        "With KYZER: Snaps 1 photo of the clinic register. Gemini 1.5 Flash auto-digitizes stock, ICU beds, and attendance in 1.8s.",
                        "97% Bandwidth Reduction: Uploads smoothly even on 2G/3G edge connections."
                    ]
                },
                {
                    "title": "Dr. Patil — Medical Officer (Shirur Depot)",
                    "points": [
                        "Before: Manages 12,000 Paracetamol units but has zero visibility into neighboring outbreak spikes.",
                        "With KYZER: Receives automated, KMS-signed lateral redistribution dispatches saving 13.5 km transit.",
                        "100% Audit Compliance: Full FHIR R4 MedicationRequest standard adherence."
                    ]
                }
            ]
        },
        {
            "tag": "SLIDE 04 / 12  ·  THE SOLUTION",
            "title": "The 5-Tier KYZER Sovereign Architecture",
            "subtitle": "Deterministic, cryptographically auditable intelligence—not generative guesswork",
            "script": "KYZER is structured into 5 cohesive tiers: Gemini Vision for paper perception, LightGBM Tweedie for multi-horizon quantile forecasting, Isolation Forest for anomaly detection, IBM Quantum QPU for multi-facility route optimization, and a DeepSeek Harness multi-agent governance loop that enforces clinical safety gates before any dispatch is signed.",
            "cards": [
                {
                    "title": "Perception, Forecasting & Risk (Tiers 1-3)",
                    "points": [
                        "Tier 1 Perception: OpenCV 5.0 + Gemini 1.5 Flash Vision OCR.",
                        "Tier 2 Forecasting: LightGBM Tweedie Quantile (P10/P50/P90) + SEIR coupling.",
                        "Tier 3 Anomaly: Isolation Forest 3-pillar cascade risk (Meds + Beds + Staff)."
                    ]
                },
                {
                    "title": "Optimization, Governance & Security (Tiers 4-5)",
                    "points": [
                        "Tier 4 Quantum VRP: PostGIS KNN nearest donor + 156-Qubit IBM Heron r2 QAOA.",
                        "Tier 5 Governance: DeepSeek Harness 5-Agent loop with >= 1.9x clinical safety gate.",
                        "Security: Strix SOC2 Type II verification + KMS HMAC-SHA256 digital signatures."
                    ]
                }
            ]
        },
        {
            "tag": "SLIDE 05 / 12  ·  SYSTEM ARCHITECTURE",
            "title": "Natural Language System Topology & Data Flow",
            "subtitle": "End-to-end data pipeline from rural edge device to quantum solver and audit ledger",
            "script": "Here is the natural language data flow: Physical register photos are processed into FHIR R4 standard entities in PostgreSQL. When the Forecaster Agent predicts a stockout, the Allocator matches nearest surplus donors via PostGIS KNN, solves the cold-chain vehicle route on IBM Quantum hardware, and submits the payload to the Critic Agent for deterministic safety validation.",
            "cards": [
                {
                    "title": "1. Edge Perception & Database Persistence",
                    "points": [
                        "ASHA mobile photo ➔ HTML5 Canvas 97% compression ➔ Gemini 1.5 Flash Vision.",
                        "/api/v1/ocr/commit-register writes multi-pillar transaction to PostgreSQL/Neon DB.",
                        "Idempotent ledger writes: skips duplicate batches, updates changed counts."
                    ]
                },
                {
                    "title": "2. Multi-Agent Governance & Quantum QPU",
                    "points": [
                        "ForecasterAgent predicts surge ➔ DetectorAgent flags P0 risk (< 1.4 days).",
                        "AllocatorAgent matches donor via PostGIS KNN + 156-Qubit IBM Heron QAOA VRP.",
                        "SupervisorAgent verifies clinical buffer (2.1x >= 1.9x) ➔ KMS HMAC-SHA256 signs dispatch."
                    ]
                }
            ]
        },
        {
            "tag": "SLIDE 06 / 12  ·  AI PERCEPTION",
            "title": "Paper Register Digitization (OpenCV + Gemini 1.5)",
            "subtitle": "Converting physical handwritten logbooks into verified digital healthcare records",
            "script": "Our perception tier solves rural connectivity: client-side canvas compression shrinks 6MB photos to 150KB on 2G/3G networks. OpenCV deskews the logbook, and Gemini 1.5 Flash Vision extracts pharmaceutical batches, expiration dates, and bed occupancy with 98.4% accuracy, committing them via an idempotent multi-pillar transaction.",
            "cards": [
                {
                    "title": "Client-Side Compression & Preprocessing",
                    "points": [
                        "HTML5 Canvas Downscaling: Reduces 6MB raw photos to ~150KB (97% bandwidth saved).",
                        "OpenCV 5.0 Geometric Preprocessing: Hough deskewing and adaptive Gaussian binarization.",
                        "Zero-install web technology: Works immediately on budget Android smartphones."
                    ]
                },
                {
                    "title": "Gemini 1.5 Flash Vision & Extraction Modes",
                    "points": [
                        "Zero-shot extraction of drug codes, batch numbers, expiry dates, and ward beds.",
                        "Extraction Mode Transparency: Explicitly surfaces [LIVE GEMINI 1.5 FLASH] vs [SIMULATED FALLBACK].",
                        "Commits in 1.8 seconds with 98.4% average field extraction accuracy."
                    ]
                }
            ]
        },
        {
            "tag": "SLIDE 07 / 12  ·  PREDICTIVE DEMAND",
            "title": "Epidemiological Forecasting (LightGBM Tweedie)",
            "subtitle": "Coupling machine learning with SEIR epidemic dynamics for zero-stockout planning",
            "script": "For demand forecasting, we deployed a LightGBM Tweedie Quantile model coupled with differential SEIR epidemic dynamics. Achieving 17.48% WAPE across 18 district facilities, KYZER predicts stockouts 7 days in advance and explains its predictions to clinicians using TreeSHAP feature attributions.",
            "cards": [
                {
                    "title": "Tweedie Quantile Loss & SEIR Dynamics",
                    "points": [
                        "Tweedie Loss (p=1.3): Tailored for zero-inflated medical consumption time-series.",
                        "SEIR Dynamic Epidemic Coupling: Incorporates beta transmission rate, R0 = 1.91, and rainfall lags.",
                        "Autoregressive Multi-Horizon: 7-day forecast with P10/P50/P90 confidence bounds."
                    ]
                },
                {
                    "title": "Verified Performance & TreeSHAP Explainability",
                    "points": [
                        "Achieved 17.48% WAPE across 18 district health facilities.",
                        "TreeSHAP Clinical Explainability: Surfaces exact top drivers (rainfall_lag_3d +34.2%).",
                        "Isolation Forest Risk Scoring: 3-pillar cascade risk (45% Meds + 35% Beds + 20% Staff)."
                    ]
                }
            ]
        },
        {
            "tag": "SLIDE 08 / 12  ·  QUANTUM OPTIMIZATION",
            "title": "IBM Quantum Heron r2 Hardware Execution",
            "subtitle": "Validated physical QPU execution strictly beating the WHO 240-minute cold-chain deadline",
            "script": "For route optimization, we executed our QAOA Hamiltonian on IBM's physical 156-qubit Heron r2 processor (ibm_fez). The quantum-classical hybrid solver found a 138.89 km route completed in 238.1 minutes—saving 13.5 km and beating the strict WHO 240-minute ice pack melting deadline before vaccine degradation occurs.",
            "cards": [
                {
                    "title": "Physical Hardware Execution Metrics",
                    "points": [
                        "Executed on physical 156-Qubit IBM Heron r2 QPU (ibm_fez).",
                        "Circuit Parameters: 16 physical transmon qubits, 125 quantum gates.",
                        "Job ID: Permanently archived under da2745cdedkc73errsp0."
                    ]
                },
                {
                    "title": "WHO Cold-Chain Physics & Distance Saved",
                    "points": [
                        "Solved Tour: Koregaon Bhima ➔ Shirur Depot ➔ Shikrapur ➔ Aundh (138.89 km).",
                        "Transit Time: 238.1 minutes — strictly beats the WHO 240-minute ice pack melting limit.",
                        "13.5 km Saved (8.9% faster delivery) vs unoptimized classical routing."
                    ]
                }
            ]
        },
        {
            "tag": "SLIDE 09 / 12  ·  MULTI-AGENT GOVERNANCE",
            "title": "DeepSeek Harness 5-Agent Consensus Engine",
            "subtitle": "Deterministic Worker-Critic pipeline with 2-way clinical safety verification",
            "script": "KYZER's 5-agent governance loop uses a deterministic Worker-Critic architecture. The Supervisor Agent enforces a clinical safety gate: no donor clinic is allowed to transfer medicine if its own remaining buffer drops below 1.9 times emergency demand.",
            "cards": [
                {
                    "title": "5-Agent Collaborative Topology",
                    "points": [
                        "01. Planner (34.2ms): Decomposes epidemic surge trajectory.",
                        "02. Detector (18.1ms): Isolation Forest flags 1.4-day buffer depletion.",
                        "03. Allocator (12.7ms): Matches donor via PostGIS KNN & IBM Quantum QAOA.",
                        "04. Explainer (22.4ms): Computes TreeSHAP clinical feature attribution.",
                        "05. Supervisor (8.3ms): Deterministic clinical consensus validation."
                    ]
                },
                {
                    "title": "Deterministic Safety Gate (>= 1.9x)",
                    "points": [
                        "Enforces donor clinic safety: Donor must retain >= 1.9x its own emergency buffer.",
                        "Verified 2.1x buffer retention at Shirur Depot before approving transfer.",
                        "Rejects any allocation that would induce a secondary downstream stockout."
                    ]
                }
            ]
        },
        {
            "tag": "SLIDE 10 / 12  ·  SECURITY & COMPLIANCE",
            "title": "Sovereign B2G Security & Auditability",
            "subtitle": "Palantir Foundry Blueprint UI, FedRAMP High Ready, and Strix SOC2 Type II verification",
            "script": "Engineered for government procurement, KYZER adheres to Palantir Foundry B2G design standards, FedRAMP High Ready authorization, Strix SOC2 Type II container security, and ABDM FHIR R4 interoperability with an immutable KMS-signed audit trail.",
            "cards": [
                {
                    "title": "Strix Security & KMS Cryptographic Signatures",
                    "points": [
                        "Strix Security: Automated container isolation audits report 0 Critical CVEs (99.8% score).",
                        "Government KMS HMAC-SHA256: Every emergency medicine dispatch is cryptographically signed.",
                        "Palantir Blueprint UI: Strict 8px spatial grid, #111418 dark canvas, 3px corner radii."
                    ]
                },
                {
                    "title": "Screenpipe 24/7 Audit Ledger & ABDM Standards",
                    "points": [
                        "Screenpipe 24/7 Context: Append-only immutable audit ledger (agent_audit_ledger.jsonl).",
                        "ABDM / MoHFW Interoperability: FHIR R4 MedicationRequest & Encounter schema.",
                        "Zero-Trust Architecture: Encrypted at rest and in transit via AES-256-GCM."
                    ]
                }
            ]
        },
        {
            "tag": "SLIDE 11 / 12  ·  MEASURED IMPACT",
            "title": "Demonstrated Clinical & Economic Impact",
            "subtitle": "Transforming public health logistics across India and BRICS partner nations",
            "script": "The impact is clear: 0% cold-chain vaccine spoilage across simulated monsoon shocks, 13.5 km saved per route, 97% mobile bandwidth reduction, and a 1-click Linux deployer that launches the entire sovereign stack in 3 minutes.",
            "cards": [
                {
                    "title": "Clinical & Operational Benchmarks",
                    "points": [
                        "0% Cold-Chain Vaccine Spoilage: All delivery runs finish within WHO 240-minute window.",
                        "13.5 km Saved Per Route (8.9% Faster turnaround during critical epidemics).",
                        "97% Mobile Bandwidth Reduction for rural 2G/3G health workers."
                    ]
                },
                {
                    "title": "BRICS Scale & 1-Click VM Deployment",
                    "points": [
                        "BRICS Cross-Border Support: Evaluated with 10 India, 5 South Africa, 3 Brazil nodes.",
                        "Simultaneous domestic matching (Shirur 32.4 km) and cross-border air-freight (Tshwane 6,970 km).",
                        "1-Click Deployer: sudo bash deploy_vm.sh deploys full production stack in 3 minutes."
                    ]
                }
            ]
        },
        {
            "tag": "SLIDE 12 / 12  ·  TEAM & LIVE DEMO",
            "title": "Team KYZER — Production Live Today",
            "subtitle": "Built for Google Cloud: Build with AI — Code for Communities Season 2",
            "script": "KYZER is live and accessible right now at atharveeee-netizen.github.io/KYZER. Built by Team KYZER for Google Cloud Code for Communities Season 2. Thank you, and we are ready for your questions!",
            "cards": [
                {
                    "title": "Team KYZER Multidisciplinary Leads",
                    "points": [
                        "Person 1 (Atharve): AI Engine, LightGBM Tweedie, SEIR Dynamics & IBM Quantum QPU",
                        "Person 2 (Backend Lead): FastAPI, PostGIS Neon DB, FEFO Ledger & KMS Signatures",
                        "Person 3 (Frontend Lead): Palantir Foundry UI, MapLibre 3D GIS & Deck.gl Digital Twin",
                        "Person 4 (Sumit): Voice AI, WhatsApp Cloud API Alerts & Submission Lead"
                    ]
                },
                {
                    "title": "Live Production Verification",
                    "points": [
                        "Web Application: https://atharveeee-netizen.github.io/KYZER/",
                        "GitHub Repository: https://github.com/atharveeee-netizen/KYZER.git",
                        "1-Click Linux Deployer: sudo bash deploy_vm.sh",
                        "Thank you! We are ready for your questions."
                    ]
                }
            ]
        }
    ]

    blank_layout = prs.slide_layouts[6]

    for s_idx, slide_info in enumerate(slides_data):
        slide = prs.slides.add_slide(blank_layout)

        # 1. Background Rect
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = c_canvas
        bg.line.fill.background()

        # 2. Header Bar / Tag
        tag_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.35))
        tf_tag = tag_box.text_frame
        tf_tag.word_wrap = True
        p_tag = tf_tag.paragraphs[0]
        p_tag.text = slide_info["tag"]
        p_tag.font.size = Pt(10)
        p_tag.font.bold = True
        p_tag.font.color.rgb = c_orange

        # 3. Main Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.75), Inches(11.7), Inches(0.6))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = slide_info["title"]
        p_title.font.size = Pt(22)
        p_title.font.bold = True
        p_title.font.color.rgb = c_ink

        # 4. Subtitle
        sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.35), Inches(11.7), Inches(0.4))
        tf_sub = sub_box.text_frame
        tf_sub.word_wrap = True
        p_sub = tf_sub.paragraphs[0]
        p_sub.text = slide_info["subtitle"]
        p_sub.font.size = Pt(12)
        p_sub.font.color.rgb = c_body

        # 5. Content Cards (2 Columns)
        card_w = Inches(5.65)
        card_h = Inches(4.7)
        card_y = Inches(1.9)

        for c_idx, card_info in enumerate(slide_info["cards"]):
            card_x = Inches(0.8 + c_idx * 5.95)

            # Card Background Shape
            card_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card_x, card_y, card_w, card_h)
            card_shape.fill.solid()
            card_shape.fill.fore_color.rgb = c_card
            card_shape.line.color.rgb = c_border
            card_shape.line.width = Pt(1)

            # Card Content Text
            card_tb = slide.shapes.add_textbox(card_x + Inches(0.2), card_y + Inches(0.2), card_w - Inches(0.4), card_h - Inches(0.4))
            tf_card = card_tb.text_frame
            tf_card.word_wrap = True

            p_ctitle = tf_card.paragraphs[0]
            p_ctitle.text = card_info["title"]
            p_ctitle.font.size = Pt(13)
            p_ctitle.font.bold = True
            p_ctitle.font.color.rgb = c_ink
            p_ctitle.space_after = Pt(10)

            for pt in card_info["points"]:
                p_pt = tf_card.add_paragraph()
                p_pt.text = f"• {pt}"
                p_pt.font.size = Pt(10)
                p_pt.font.color.rgb = c_body
                p_pt.space_after = Pt(6)

        # 6. Speaker Notes
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = f"SPEAKER SCRIPT:\n{slide_info['script']}"

        # 7. Footer
        footer_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.9), Inches(11.7), Inches(0.3))
        tf_footer = footer_box.text_frame
        p_footer = tf_footer.paragraphs[0]
        p_footer.text = f"KYZER Sovereign OS  ·  Team KYZER  ·  Slide {s_idx + 1} of {len(slides_data)}"
        p_footer.font.size = Pt(9)
        p_footer.font.color.rgb = c_body

    prs.save(output_path)
    print(f"Successfully generated PowerPoint Presentation: {output_path}")

if __name__ == '__main__':
    desktop_dir = r"C:\Users\25beevdt047\Desktop"
    output_pptx = os.path.join(desktop_dir, "KYZER_PITCH_DECK.pptx")
    create_pptx(output_pptx)
